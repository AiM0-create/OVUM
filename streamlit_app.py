import streamlit as st
import geopandas as gpd
import pandas as pd
import osmnx as ox
from shapely.geometry import Point
from pptx import Presentation
from pptx.util import Inches
import zipfile
import tempfile

st.set_page_config(page_title="Catchment Analysis", layout="wide")

st.title("Catchment Analysis Dashboard")

# Sidebar: File uploads
st.sidebar.header("Upload Data")
csv_file = st.sidebar.file_uploader("Upload clinic CSV", type=["csv"])
shp_zip = st.sidebar.file_uploader("Upload catchment shapefile ZIP", type=["zip"])

def load_shapefile(zip_file):
    with tempfile.TemporaryDirectory() as tmpdir:
        z = zipfile.ZipFile(zip_file)
        z.extractall(tmpdir)
        # find .shp
        for fname in z.namelist():
            if fname.endswith('.shp'):
                shp_path = f"{tmpdir}/{fname}"
                return gpd.read_file(shp_path)
    return None

@st.cache_data
def load_clinics(csv_bytes):
    df = pd.read_csv(csv_bytes)
    gdf = gpd.GeoDataFrame(df,
                           geometry=gpd.points_from_xy(df['long'], df['lat']),
                           crs="EPSG:4326")
    gdf_proj = gdf.to_crs(epsg=3857)
    for meters in [3000, 5000, 6000]:
        label = meters//1000
        gdf_proj[f'buffer_{label}km'] = gdf_proj.geometry.buffer(meters)
    return gdf_proj

@st.cache_data
def extract_osm_metrics(clinics_gdf, km_label):
    records = []
    for _, row in clinics_gdf.iterrows():
        poly = row[f'buffer_{km_label}km'].to_crs(epsg=4326)
        bldgs = ox.geometries.geometries_from_polygon(poly, {'building': True})
        comps = ox.geometries.geometries_from_polygon(poly, {'amenity':'hospital','healthcare':'maternity'})
        G = ox.graph_from_polygon(poly, network_type='drive')
        length_km = sum(d['length'] for _,_,d in G.edges(data=True))/1000
        area_km2 = poly.to_crs(epsg=3857).area/1e6
        schools = ox.geometries.geometries_from_polygon(poly, {'amenity':['kindergarten','school']})
        pts = ox.geometries.geometries_from_polygon(poly, {'public_transport':'station'})
        records.append({
            'clinic': row['name'],
            'radius_km': km_label,
            'buildings': len(bldgs),
            'competitors': len(comps),
            'road_density': length_km/area_km2,
            'playschools': len(schools),
            'public_transport': len(pts)
        })
    return pd.DataFrame(records)

@st.cache_data
def compute_scores(df):
    weights = {'buildings':0.2,'competitors':-0.3,'road_density':0.2,'playschools':0.1,'public_transport':0.2}
    df2 = df.copy()
    for k,w in weights.items():
        norm=(df2[k]-df2[k].min())/(df2[k].max()-df2[k].min())
        df2[f'{k}_score']=norm*w
    df2['total_score']=df2[[f'{k}_score' for k in weights]].sum(axis=1)
    df2['rank']=df2['total_score'].rank(ascending=False)
    return df2

if csv_file and shp_zip:
    st.success("Data loaded successfully.")
    clinics = load_clinics(csv_file)
    gdf_catch = load_shapefile(shp_zip)

    st.subheader("Clinic Locations & Catchments")
    st.map(pd.DataFrame({'lat':clinics.to_crs(epsg=4326).geometry.y,'lon':clinics.to_crs(epsg=4326).geometry.x}))

    # Extract metrics
    all_metrics = pd.concat([extract_osm_metrics(clinics, r) for r in [3,5,6]], ignore_index=True)
    scores = compute_scores(all_metrics)

    st.subheader("Catchment Metrics & Scores")
    st.dataframe(scores)

    # Download PowerPoint report
    def create_ppt(scores_df):
        prs = Presentation()
        for clinic in scores_df['clinic'].unique():
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = clinic
            subset = scores_df[scores_df['clinic']==clinic]
            rows,cols=subset.shape[0]+1,4
            table=slide.shapes.add_table(rows,cols,Inches(0.5),Inches(1.5),Inches(9),Inches(0.8+0.2*rows)).table
            hdr=['Radius (km)','Score','Competitors','Road Density']
            for i,h in enumerate(hdr): table.cell(0,i).text=h
            for i,row in enumerate(subset.itertuples(),start=1):
                table.cell(i,0).text=str(row.radius_km)
                table.cell(i,1).text=f"{row.total_score:.2f}"
                table.cell(i,2).text=str(row.competitors)
                table.cell(i,3).text=f"{row.road_density:.2f}"
        return prs

    if st.button("Generate PPT Report"):
        prs = create_ppt(scores)
        buf = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
        prs.save(buf.name)
        st.success("Report generated!")
        st.download_button("Download PPTX", buf.name, file_name="catchment_report.pptx")
else:
    st.info("Please upload both the clinic CSV and shapefile ZIP to proceed.")
